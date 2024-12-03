from fastapi import APIRouter, File, UploadFile, HTTPException
from fastapi.responses import JSONResponse
from io import BytesIO
import google.generativeai as genai
from utilities.env import Env
from typing import List, Dict
import typing
import PyPDF2
import json
import re
import tqdm
import pptx

cards_router = APIRouter()
env = Env()


class StudyMaterialGenerator:
    def __init__(self, api_key: str, model_name: str = "gemini-pro"):
        if not api_key:
            raise ValueError("Gemini API key is required")

        genai.configure(api_key=api_key)

        try:
            self.model = genai.GenerativeModel(model_name)
        except Exception as e:
            raise RuntimeError(f"Failed to initialize Gemini model: {e}")

    def _safe_generate(self, prompt: str, max_tokens: int = 2048) -> str:
        safety_settings = [
            {"category": "HARM_CATEGORY_HARASSMENT", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_HATE_SPEECH", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_SEXUALLY_EXPLICIT", "threshold": "BLOCK_NONE"},
            {"category": "HARM_CATEGORY_DANGEROUS_CONTENT", "threshold": "BLOCK_NONE"},
        ]

        generation_config = {
            "max_output_tokens": max_tokens,
            "temperature": 0.7,
            "top_p": 0.9,
        }

        try:
            response = self.model.generate_content(
                prompt,
                generation_config=generation_config,
                safety_settings=safety_settings,
            )
            return response.text
        except Exception as e:
            print(f"Error generating content: {e}")
            return ""

    def _split_text_into_chunks(self, text: str, max_length: int = 2000) -> List[str]:
        chunks = []
        while len(text) > max_length:
            split_point = text.rfind(" ", max_length - 100, max_length)
            if split_point == -1:
                split_point = max_length
            chunks.append(text[:split_point])
            text = text[split_point:].strip()
        if text:
            chunks.append(text)
        return chunks

    def generate_flashcards(self, context: str) -> List[Dict[str, str]]:
        prompt = f"""Create a JSON array of flashcards from the following content.
        Each flashcard should be a dictionary with 'question' and 'answer' keys.

        Content:
        {context}

        Output format:
        [
            {{"question": "...", "answer": "..."}}
        ]"""

        print(f"Flashcard prompt: {prompt}")
        response = self._safe_generate(prompt)
        print(f"Flashcard response: {response}")
        json_match = re.search(r"\[.*\]", response, re.DOTALL)
        if json_match:
            try:
                return json.loads(json_match.group(0))
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}")
                print(f"Problematic response: {json_match.group(0)}")
                return []
        return []

    def generate_quiz(self, context: str) -> List[Dict[str, str]]:
        prompt = f"""Create a JSON array of multiple-choice quiz questions from the following content.
        Each question should be a dictionary with 'question', 'options', and 'correct_answer' keys.

        Content:
        {context}

        Output format:
        [
            {{
                "question": "...",
                "options": ["A) ...", "B) ...", "C) ...", "D) ..."],
                "correct_answer": "A"
            }}
        ]"""

        print(f"Quiz prompt: {prompt}")
        response = self._safe_generate(prompt)
        print(f"Quiz response: {response}")
        json_match = re.search(r"\[.*\]", response, re.DOTALL)
        if json_match:
            try:
                return json.loads(json_match.group(0))
            except json.JSONDecodeError as e:
                print(f"JSONDecodeError: {e}")
                print(f"Problematic response: {json_match.group(0)}")
                return []
        return []

    def generate_notes(self, context: str) -> str:
        prompt = f"""Generate comprehensive, structured study notes from the following content.
        Use markdown formatting with headings, bullet points, and clear sections.

        Content:
        {context}"""

        print(f"Notes prompt: {prompt}")
        response = self._safe_generate(prompt)
        print(f"Notes response: {response}")
        return response

    def process_large_text(self, text: str) -> typing.Dict[str, typing.Any]:
        chunks = self._split_text_into_chunks(text)
        print(f"Text chunks: {chunks}")
        all_flashcards = []
        all_quiz = []
        all_notes = []

        for chunk in chunks:
            flashcards = self.generate_flashcards(chunk)
            quiz = self.generate_quiz(chunk)
            notes = self.generate_notes(chunk)

            all_flashcards.extend(flashcards)
            all_quiz.extend(quiz)
            all_notes.append(notes)

        return {
            "flashcards": all_flashcards,
            "quiz": all_quiz,
            "notes": "\n".join(all_notes),
        }


def extract_text_from_pdf(pdf_stream: BytesIO) -> str:
    """Extract text from a PDF file using a BytesIO stream."""
    text = ""
    reader = PyPDF2.PdfReader(pdf_stream)
    for page in tqdm.tqdm(reader.pages, desc="Reading PDF pages"):
        text += page.extract_text() or ""
    print(f"Extracted text from PDF: {text}")
    return text.strip()


def extract_text_from_pptx(pptx_stream: BytesIO) -> str:
    """Extract text from a PPTX file using a BytesIO stream."""
    text = ""
    presentation = pptx.Presentation(pptx_stream)
    for slide in tqdm.tqdm(presentation.slides, desc="Reading PPTX slides"):
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    print(f"Extracted text from PPTX: {text}")
    return text.strip()


@cards_router.post("/upload/")
async def upload_file(file: UploadFile = File(...)):
    api_key = env.GEMINI_API_KEY
    if not api_key:
        raise HTTPException(status_code=400, detail="Gemini API key is required")

    generator = StudyMaterialGenerator(api_key=api_key)
    print(f"StudyMaterialGenerator initialized with API key: {api_key}")

    file_extension = file.filename.split(".")[-1].lower()
    file_content = await file.read()

    if file_extension == "pdf":
        text = extract_text_from_pdf(BytesIO(file_content))
    elif file_extension == "pptx":
        text = extract_text_from_pptx(BytesIO(file_content))
    else:
        raise HTTPException(status_code=400, detail="Unsupported file type")

    result = generator.process_large_text(text)
    return JSONResponse(content=result)


@cards_router.post("/text/")
async def upload(query: str):
    api_key = env.GEMINI_API_KEY
    if not api_key:
        raise HTTPException(status_code=400, detail="Gemini API key is required")

    generator = StudyMaterialGenerator(api_key=api_key)
    print(f"StudyMaterialGenerator initialized with API key: {api_key}")

    if query:
        text = query
    else:
        raise HTTPException(status_code=400, detail="Context must be provided")

    result = generator.process_large_text(text)
    return JSONResponse(content=result)
